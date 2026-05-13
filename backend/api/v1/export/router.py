from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse
from typing import Dict, Any
from pydantic import BaseModel
import os
from datetime import datetime

router = APIRouter()

EXPORT_DIR = "exports"
os.makedirs(EXPORT_DIR, exist_ok=True)


class ExportRequest(BaseModel):
    content: str
    format: str
    filename: str = None


@router.post("/markdown")
async def export_markdown(request: ExportRequest):
    filename = request.filename or f"export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md"
    filepath = os.path.join(EXPORT_DIR, filename)

    try:
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(request.content)

        return {
            "status": "success",
            "format": "markdown",
            "filename": filename,
            "filepath": filepath
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/docx")
async def export_docx(request: ExportRequest):
    try:
        from docx import Document
        from docx.shared import Pt

        doc = Document()
        filename = request.filename or f"export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        filepath = os.path.join(EXPORT_DIR, filename)

        for paragraph in request.content.split("\n"):
            if paragraph.strip():
                p = doc.add_paragraph()
                run = p.add_run(paragraph.strip())
                if paragraph.startswith("# "):
                    run.font.size = Pt(18)
                    run.font.bold = True
                elif paragraph.startswith("## "):
                    run.font.size = Pt(16)
                    run.font.bold = True
                elif paragraph.startswith("### "):
                    run.font.size = Pt(14)
                    run.font.bold = True

        doc.save(filepath)

        return {
            "status": "success",
            "format": "docx",
            "filename": filename,
            "filepath": filepath
        }
    except ImportError:
        raise HTTPException(status_code=500, detail="python-docx not installed")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/pptx")
async def export_pptx(request: ExportRequest):
    try:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation()
        filename = request.filename or f"export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(EXPORT_DIR, filename)

        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        lines = request.content.split("\n")
        if lines:
            title.text = lines[0].replace("# ", "").strip()
        if len(lines) > 1:
            subtitle.text = lines[1].strip()

        for i in range(2, len(lines), 5):
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = "内容"

            tf = body_shape.text_frame
            for j in range(i, min(i + 5, len(lines))):
                if lines[j].strip():
                    p = tf.add_paragraph()
                    p.text = lines[j].strip()
                    p.level = 0

        prs.save(filepath)

        return {
            "status": "success",
            "format": "pptx",
            "filename": filename,
            "filepath": filepath
        }
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx not installed")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/download/{filename}")
async def download_file(filename: str):
    filepath = os.path.join(EXPORT_DIR, filename)
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail=f"File {filename} not found")

    return FileResponse(
        path=filepath,
        filename=filename,
        media_type="application/octet-stream"
    )


@router.get("/list")
async def list_exports():
    files = []
    for filename in os.listdir(EXPORT_DIR):
        filepath = os.path.join(EXPORT_DIR, filename)
        if os.path.isfile(filepath):
            stat = os.stat(filepath)
            files.append({
                "filename": filename,
                "size": stat.st_size,
                "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat()
            })
    return {"exports": files, "count": len(files)}
