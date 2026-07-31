"""Markdown 转换器 — V3 实装版（行业标准工具）

使用 pandoc / marp / markmap 三个行业标准工具生成高质量导出：
- md → docx: pypandoc-binary（pandoc 是文档转换行业标准，保留完整 markdown 特性）
- md → pptx: marp-cli（最佳 markdown → slides 工具，支持主题/代码高亮/数学公式）
- md → mindmap: markmap-cli（交互式思维导图 HTML，可折叠/缩放）

每个主转换函数都有 fallback 到纯 Python 实现（当工具不可用时）。

依赖：
    必需：pypandoc-binary (pip install)
    必需：@marp-team/marp-cli (npm install -g)
    必需：markmap-cli (npm install -g)
    Fallback：python-docx, python-pptx, pyecharts (pip install)

设计原则：
    1. 优先使用行业标准工具（pandoc/marp/markmap）保证输出质量
    2. 工具不可用时降级到纯 Python 实现，保证功能可用
    3. 统一错误处理和日志
"""
from __future__ import annotations

import logging
import os
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from typing import List, Optional, Tuple

logger = logging.getLogger(__name__)


# =============================================================================
# Markdown → DOCX (pandoc)
# =============================================================================

def convert_to_docx(md_text: str, output_path: str) -> str:
    """用 pandoc 把 Markdown 转换为 docx（行业标准质量）

    pandoc 保留完整 markdown 特性：
    - 标题层级（自动映射到 Word 标题样式）
    - 嵌套列表、有序/无序列表
    - 表格（含对齐）
    - 代码块（含语法高亮信息）
    - 引用块（嵌套）
    - 链接、图片、脚注
    - 分隔线

    Args:
        md_text: Markdown 文本
        output_path: 输出 .docx 文件路径

    Returns:
        output_path

    Fallback:
        pandoc 不可用时降级到 python-docx 纯 Python 实现
    """
    try:
        import pypandoc
        # pypandoc-binary 自带 pandoc 二进制，无需额外安装
        pypandoc.convert_text(
            md_text,
            'docx',
            format='md',
            outputfile=output_path,
            extra_args=['--standalone'],
        )
        logger.info(f"DOCX generated via pandoc: {output_path}")
        return output_path
    except ImportError:
        logger.warning("pypandoc not available, falling back to python-docx")
        return _fallback_convert_to_docx(md_text, output_path)
    except Exception as e:
        logger.warning(f"pandoc conversion failed ({e}), falling back to python-docx")
        return _fallback_convert_to_docx(md_text, output_path)


# =============================================================================
# Markdown → PPTX (marp)
# =============================================================================

def convert_to_pptx(md_text: str, output_path: str) -> str:
    """用 marp 把 Markdown 转换为 pptx（最佳 md→slides 质量）

    marp 特性：
    - 智能分页（--- 或 H2 自动分页）
    - 代码语法高亮
    - 数学公式（KaTeX）
    - 主题支持（default/gaia/uncover）
    - 页码、页眉页脚
    - 图片自适应

    转换策略：
    1. 把普通 markdown 转成 marp 格式（加 frontmatter + H2 前加分页符）
    2. 调用 marp CLI 生成 pptx

    Args:
        md_text: Markdown 文本
        output_path: 输出 .pptx 文件路径

    Returns:
        output_path

    Fallback:
        marp 不可用时降级到 python-pptx 纯 Python 实现
    """
    marp_path = shutil.which('marp')
    if not marp_path:
        logger.warning("marp CLI not found in PATH, falling back to python-pptx")
        return _fallback_convert_to_pptx(md_text, output_path)

    # 把 markdown 转成 marp 格式
    marp_md = _to_marp_format(md_text)

    # 写入临时 md 文件
    with tempfile.NamedTemporaryFile(
        mode='w', suffix='.md', delete=False, encoding='utf-8'
    ) as f:
        f.write(marp_md)
        temp_md = f.name

    try:
        result = subprocess.run(
            [
                marp_path,
                '--pptx',
                '--allow-local-files',
                '-o', output_path,
                temp_md,
            ],
            capture_output=True,
            text=True,
            timeout=120,  # 2 分钟超时
        )
        if result.returncode != 0:
            logger.warning(
                f"marp failed (exit {result.returncode}): {result.stderr}, "
                f"falling back to python-pptx"
            )
            return _fallback_convert_to_pptx(md_text, output_path)

        logger.info(f"PPTX generated via marp: {output_path}")
        return output_path
    except subprocess.TimeoutExpired:
        logger.warning("marp timed out (120s), falling back to python-pptx")
        return _fallback_convert_to_pptx(md_text, output_path)
    finally:
        if os.path.exists(temp_md):
            os.unlink(temp_md)


def _to_marp_format(md_text: str) -> str:
    """把普通 markdown 转成 marp 兼容格式

    - 添加 marp frontmatter（启用 marp + 主题 + 分页）
    - 移除原始 markdown 中的 --- 分隔线（避免被 marp 识别为分页符产生空白页）
    - 在每个 H2 前插入分页符 (---)
    - H1 作为标题页

    修复：LLM 生成的 markdown 常在段落间使用 --- 作为分隔线，
    这些分隔线会被 marp 识别为分页符，导致生成大量空白页。
    因此需要先清除原始 --- 分隔线，再统一在 H2 前插入分页符。
    """
    frontmatter = """---
marp: true
theme: default
paginate: true
size: 16:9
backgroundColor: #ffffff
style: |
  section {
    font-family: "Microsoft YaHei", "PingFang SC", sans-serif;
    color: #333;
  }
  h1 {
    color: #6366f1;
    border-bottom: 2px solid #6366f1;
    padding-bottom: 8px;
  }
  h2 {
    color: #8b5cf6;
  }
  code {
    background: #f5f5f5;
    padding: 2px 6px;
    border-radius: 3px;
  }
  pre {
    background: #1e1e1e;
    color: #d4d4d4;
    border-radius: 6px;
  }
  blockquote {
    border-left: 4px solid #8b5cf6;
    background: rgba(139, 92, 246, 0.05);
    padding: 8px 16px;
  }
  table {
    border-collapse: collapse;
  }
  th {
    background: #6366f1;
    color: white;
  }
  th, td {
    border: 1px solid #ddd;
    padding: 8px 12px;
  }
---

"""
    lines = md_text.split('\n')
    result = []

    for line in lines:
        stripped = line.strip()

        # 跳过原始 markdown 中的 --- 分隔线（thematic break）
        # 这些分隔线会被 marp 识别为分页符，导致空白页
        # 同时跳过 *** 和 ___ 分隔线（markdown 规范中的其他分隔线形式）
        if re.match(r'^(-{3,}|\*{3,}|_{3,})$', stripped):
            continue

        # H2 前加分页符（跳过 H3+，它们属于同一页）
        if line.startswith('## ') and not line.startswith('### '):
            # 确保分页符前有且仅有一个空行
            if result and result[-1] != '':
                result.append('')
            result.append('---')
            result.append('')

        result.append(line)

    # 清理末尾多余空行
    while result and result[-1] == '':
        result.pop()

    return frontmatter + '\n'.join(result)


# =============================================================================
# Markdown → 思维导图 HTML (markmap)
# =============================================================================

def convert_to_mindmap(md_text: str, output_path: str) -> str:
    """用 markmap 把 Markdown 转换为交互式思维导图 HTML

    markmap 特性：
    - 基于 markdown 标题层级自动构建树
    - 交互式：可折叠/展开节点
    - 可缩放、可拖拽
    - 支持代码高亮、数学公式
    - 自动布局优化

    Args:
        md_text: Markdown 文本
        output_path: 输出 .html 文件路径

    Returns:
        output_path

    Fallback:
        markmap 不可用时降级到 pyecharts Tree 实现
    """
    markmap_path = shutil.which('markmap')
    if not markmap_path:
        logger.warning("markmap CLI not found in PATH, falling back to pyecharts")
        return _fallback_convert_to_mindmap(md_text, output_path)

    # 写入临时 md 文件
    with tempfile.NamedTemporaryFile(
        mode='w', suffix='.md', delete=False, encoding='utf-8'
    ) as f:
        f.write(md_text)
        temp_md = f.name

    try:
        result = subprocess.run(
            [
                markmap_path,
                temp_md,
                '-o', output_path,
            ],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode != 0:
            logger.warning(
                f"markmap failed (exit {result.returncode}): {result.stderr}, "
                f"falling back to pyecharts"
            )
            return _fallback_convert_to_mindmap(md_text, output_path)

        logger.info(f"Mindmap HTML generated via markmap: {output_path}")
        return output_path
    except subprocess.TimeoutExpired:
        logger.warning("markmap timed out (60s), falling back to pyecharts")
        return _fallback_convert_to_mindmap(md_text, output_path)
    finally:
        if os.path.exists(temp_md):
            os.unlink(temp_md)


# =============================================================================
# Fallback：纯 Python 实现（当 pandoc/marp/markmap 不可用时）
# =============================================================================

@dataclass
class _MdNode:
    """Fallback 解析用的 Markdown 节点"""
    type: str
    level: int = 0
    text: str = ''
    items: List[str] = field(default_factory=list)
    ordered: bool = False
    rows: List[List[str]] = field(default_factory=list)


def _parse_markdown(md_text: str) -> List[_MdNode]:
    """纯 Python 行级 Markdown 解析（fallback 用）"""
    lines = md_text.split('\n')
    nodes: List[_MdNode] = []
    i = 0
    n = len(lines)
    list_buffer: List[str] = []
    list_ordered = False

    def flush_list():
        nonlocal list_buffer, list_ordered
        if list_buffer:
            nodes.append(_MdNode(type='list', ordered=list_ordered, items=list_buffer))
            list_buffer = []
            list_ordered = False

    while i < n:
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            flush_list()
            i += 1
            continue

        if stripped.startswith('```'):
            flush_list()
            code_lines = []
            i += 1
            while i < n and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            nodes.append(_MdNode(type='code_block', text='\n'.join(code_lines)))
            i += 1
            continue

        m = re.match(r'^(#{1,6})\s+(.+?)\s*#*$', stripped)
        if m:
            flush_list()
            nodes.append(_MdNode(
                type='heading',
                level=len(m.group(1)),
                text=_strip_inline_md(m.group(2)),
            ))
            i += 1
            continue

        if re.match(r'^(-{3,}|\*{3,}|_{3,})$', stripped):
            flush_list()
            nodes.append(_MdNode(type='thematic_break'))
            i += 1
            continue

        if '|' in stripped and stripped.startswith('|') and i + 1 < n:
            next_line = lines[i + 1].strip()
            if re.match(r'^\|[\s\-:|]+\|?$', next_line):
                flush_list()
                rows = [_parse_table_row(stripped)]
                i += 2
                while i < n and lines[i].strip().startswith('|'):
                    rows.append(_parse_table_row(lines[i].strip()))
                    i += 1
                nodes.append(_MdNode(type='table', rows=rows))
                continue

        if stripped.startswith('>'):
            flush_list()
            quote_text = stripped.lstrip('>').strip()
            i += 1
            while i < n and lines[i].strip().startswith('>'):
                quote_text += '\n' + lines[i].strip().lstrip('>').strip()
                i += 1
            nodes.append(_MdNode(type='blockquote', text=_strip_inline_md(quote_text)))
            continue

        m = re.match(r'^[-*+]\s+(.+)$', stripped)
        if m:
            if not list_buffer:
                list_ordered = False
            list_buffer.append(_strip_inline_md(m.group(1)))
            i += 1
            continue

        m = re.match(r'^\d+\.\s+(.+)$', stripped)
        if m:
            if not list_buffer:
                list_ordered = True
            list_buffer.append(_strip_inline_md(m.group(1)))
            i += 1
            continue

        flush_list()
        para_lines = [stripped]
        i += 1
        while i < n:
            ns = lines[i].strip()
            if (not ns or ns.startswith('#') or ns.startswith('```')
                or ns.startswith('|') or ns.startswith('>')
                or ns.startswith('---') or re.match(r'^[-*+]\s', ns)
                or re.match(r'^\d+\.\s', ns)):
                break
            para_lines.append(ns)
            i += 1
        nodes.append(_MdNode(type='paragraph', text=_strip_inline_md(' '.join(para_lines))))

    flush_list()
    return nodes


def _parse_table_row(line: str) -> List[str]:
    line = line.strip()
    if line.startswith('|'):
        line = line[1:]
    if line.endswith('|'):
        line = line[:-1]
    return [c.strip() for c in line.split('|')]


def _strip_inline_md(text: str) -> str:
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'__([^_]+)__', r'\1', text)
    text = re.sub(r'\*([^*]+)\*', r'\1', text)
    text = re.sub(r'_([^_]+)_', r'\1', text)
    text = re.sub(r'`([^`]+)`', r'\1', text)
    text = re.sub(r'!\[([^\]]*)\]\([^)]+\)', r'\1', text)
    return text.strip()


def _fallback_convert_to_docx(md_text: str, output_path: str) -> str:
    """Fallback: 用 python-docx 生成 docx"""
    from docx import Document
    from docx.shared import Pt, RGBColor

    doc = Document()
    style = doc.styles['Normal']
    style.font.name = '微软雅黑'
    style.font.size = Pt(11)

    for node in _parse_markdown(md_text):
        if node.type == 'heading':
            levels = {1: ('Heading 1', 18), 2: ('Heading 2', 16), 3: ('Heading 3', 14),
                      4: ('Heading 4', 13), 5: ('Heading 5', 12), 6: ('Heading 6', 11)}
            sn, sz = levels.get(node.level, ('Heading 6', 11))
            p = doc.add_paragraph(style=sn)
            p.add_run(node.text).font.size = Pt(sz)
        elif node.type == 'paragraph':
            doc.add_paragraph(node.text)
        elif node.type == 'list':
            for i, item in enumerate(node.items):
                prefix = f"{i+1}. " if node.ordered else "• "
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Pt(20)
                p.add_run(prefix + item)
        elif node.type == 'table' and node.rows:
            cols = max(len(r) for r in node.rows)
            table = doc.add_table(rows=len(node.rows), cols=cols)
            table.style = 'Light Grid Accent 1'
            for i, row in enumerate(node.rows):
                for j, cell in enumerate(row):
                    if j < cols:
                        table.rows[i].cells[j].text = cell
        elif node.type == 'code_block':
            p = doc.add_paragraph()
            run = p.add_run(node.text)
            run.font.name = 'Consolas'
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        elif node.type == 'blockquote':
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Pt(30)
            run = p.add_run(node.text)
            run.font.italic = True
            run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    doc.save(output_path)
    logger.info(f"DOCX generated via python-docx (fallback): {output_path}")
    return output_path


def _fallback_convert_to_pptx(md_text: str, output_path: str) -> str:
    """Fallback: 用 python-pptx 生成 pptx（按 H2 分页）"""
    from pptx import Presentation
    from pptx.util import Pt, Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    nodes = _parse_markdown(md_text)
    title_text = 'AgentMatrix 生成内容'
    for node in nodes:
        if node.type == 'heading' and node.level == 1:
            title_text = node.text
            break

    # 标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = '由 AgentMatrix 多智能体协同平台生成'

    current_title = None
    bullets: List[Tuple[str, int]] = []

    def flush():
        nonlocal current_title, bullets
        if current_title or bullets:
            s = prs.slides.add_slide(prs.slide_layouts[1])
            s.shapes.title.text = current_title or '内容'
            tf = s.shapes.placeholders[1].text_frame
            tf.word_wrap = True
            for i, (text, lvl) in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.level = max(0, lvl - 1)
                for r in p.runs:
                    r.font.size = Pt(18 - lvl * 2)
        bullets = []

    for node in nodes:
        if node.type == 'heading':
            if node.level == 2:
                flush()
                current_title = node.text
            elif node.level >= 3:
                if len(bullets) >= 6:
                    flush()
                bullets.append((node.text, node.level - 2))
        elif node.type == 'list':
            for item in node.items:
                if len(bullets) >= 6:
                    flush()
                bullets.append((item, 1))
        elif node.type == 'paragraph':
            if len(bullets) >= 6:
                flush()
            bullets.append((node.text, 1))
    flush()

    prs.save(output_path)
    logger.info(f"PPTX generated via python-pptx (fallback): {output_path}")
    return output_path


def _fallback_convert_to_mindmap(md_text: str, output_path: str) -> str:
    """Fallback: 用 pyecharts Tree 生成思维导图 HTML"""
    nodes = _parse_markdown(md_text)
    root = {'name': 'AgentMatrix 内容', 'children': []}
    stack = [root]

    def add_to_level(level: int, name: str):
        while len(stack) > 1 and len(stack) - 1 >= level:
            stack.pop()
        new_node = {'name': name, 'children': []}
        stack[-1]['children'].append(new_node)
        stack.append(new_node)

    for node in nodes:
        if node.type == 'heading':
            if node.level == 1 and not root['children']:
                root['name'] = node.text
                continue
            add_to_level(node.level, node.text)
        elif node.type == 'list':
            for item in node.items:
                stack[-1]['children'].append({'name': item[:50], 'children': []})
        elif node.type == 'paragraph':
            text = node.text[:80] + ('...' if len(node.text) > 80 else '')
            stack[-1]['children'].append({'name': text, 'children': []})

    if not root['children']:
        root['children'].append({'name': '（无内容）', 'children': []})

    from pyecharts import options as opts
    from pyecharts.charts import Tree

    tree = (
        Tree(init_opts=opts.InitOpts(width="100%", height="800px"))
        .add(
            "",
            [root],
            orient="LR",
            collapse_interval=2,
            label_opts=opts.LabelOpts(
                position="top",
                horizontal_align="right",
                vertical_align="middle",
                font_size=12,
            ),
            symbol_size=8,
        )
        .set_global_opts(
            title_opts=opts.TitleOpts(
                title=root['name'],
                subtitle="由 AgentMatrix 生成 · 基于 pyecharts Tree (fallback)",
                pos_left="center",
            ),
            tooltip_opts=opts.TooltipOpts(trigger="item", formatter="{b}"),
        )
    )
    tree.render(output_path)
    logger.info(f"Mindmap HTML generated via pyecharts (fallback): {output_path}")
    return output_path
